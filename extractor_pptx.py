from pptx import Presentation
from pathlib import Path

def extract_from_pptx(pptx_path: str):
    prs = Presentation(pptx_path)
    slides = []
    thumbs = Path("out/thumbs")
    thumbs.mkdir(parents=True, exist_ok=True)

    for i, slide in enumerate(prs.slides, start=1):
        texts, imgs = [], []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = shape.text.strip()
                if t:
                    texts.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        c = (cell.text or "").strip()
                        if c:
                            texts.append(c)

        if getattr(slide, "notes_slide", None) and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text
            if note:
                texts.append(f"[Notes] {note.strip()}")

        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                img = shape.image
                ext = img.ext or "png"
                outp = thumbs / f"slide{i}_{len(imgs)+1}.{ext}"
                outp.write_bytes(img.blob)
                imgs.append(str(outp))

        slides.append({"slide": i, "text": "\n".join(texts), "images": imgs})
    return slides
